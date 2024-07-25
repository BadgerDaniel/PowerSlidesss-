import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import numpy as np
from io import BytesIO

# Function to add footnote to slides
def add_footnote(slide, prs, text="Your footnote text here"):
    slide_height = prs.slide_height
    slide_width = prs.slide_width
    footnote_height = Emu(500000)
    footnote_top = slide_height - footnote_height
    footnote_left = Emu(0)
    footnote_width = slide_width

    textbox = slide.shapes.add_textbox(footnote_left, footnote_top, footnote_width, footnote_height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

# Function to create presentation
def create_presentation(df, footnote_text, year):
    prs = Presentation()
    # some constants
    font_size_header = Pt(12)
    font_size_body = Pt(9.5)
    table_left = Inches(1.1)
    table_top = Inches(.70)
    table_width = Inches(8.5)
    row_height = Emu(140000)
    rows_per_table = 22

    for group in df['Section'].unique():
        grouped_df = df[df['Section'] == group]
        num_tables = len(grouped_df) // rows_per_table + (len(grouped_df) % rows_per_table > 0)

        for i in range(num_tables):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            shift_left = Emu(-3000000)
            shift_up = Emu(200000)

            textbox1 = slide.shapes.add_textbox(Emu(507245) - shift_left, Emu(463902) - shift_up, Emu(10509370), Emu(436017))
            tf1 = textbox1.text_frame
            tf1.text = group
            p = tf1.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(22)
            run.font.name = 'Avenir Next LT Pro'

            shift_left = Emu(550000)
            shift_up = Emu(300000)
            
            textbox2 = slide.shapes.add_textbox(Emu(492980) - shift_left, Emu(260671) - shift_up, Emu(4383819), Emu(202944))
            tf2 = textbox2.text_frame
            tf2.text = f'{year} About Site Traffic Report'
            p = tf2.paragraphs[0]
            run = p.runs[0]
            run.font.size = Pt(19)
            run.font.name = 'Avenir Next LT Pro'
            run.font.color.rgb = RGBColor(0xD1, 0x41, 0x59)

            add_footnote(slide, prs, text=footnote_text)

            start_idx = i * rows_per_table
            end_idx = min(start_idx + rows_per_table, len(grouped_df))
            table_height = row_height * (end_idx - start_idx + 1)

            table = slide.shapes.add_table(rows=end_idx - start_idx + 1, cols=2, 
                                           left=table_left, top=table_top, 
                                           width=table_width, height=table_height).table

            first_col_width = table_width * 0.80
            second_col_width = table_width * 0.20
            table.columns[0].width = int(first_col_width)
            table.columns[1].width = int(second_col_width)

            for row_idx, row in enumerate(table.rows):
                row.height = row_height
                for col_idx, cell in enumerate(row.cells):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx > 0 else RGBColor(0, 0, 0)
                    
                    text_frame = cell.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()

                    if row_idx == 0:
                        run.text = ['Pages', 'Visits'][col_idx]
                        run.font.size = font_size_header
                        p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                    else:
                        if col_idx == 0:
                            try:
                                cell_value = grouped_df.iloc[row_idx - 1 + start_idx]['Row Labels']
                            except:
                                cell_value = grouped_df.iloc[row_idx - 1 + start_idx]['Page']
                        else:
                            try:
                                cell_value = grouped_df.iloc[row_idx - 1 + start_idx]['Adobe Visits']
                            except:
                                cell_value = grouped_df.iloc[row_idx - 1 + start_idx]['Visits']

                        if isinstance(cell_value, (int, float, np.integer)):
                            run.font.size = font_size_body
                            run.font.bold = True
                            p.alignment = PP_ALIGN.CENTER if col_idx == 1 else PP_ALIGN.LEFT
                            run.text = f"{cell_value:,.0f}"
                        else:
                            run.font.size = font_size_body
                            p.alignment = PP_ALIGN.CENTER if col_idx == 1 else PP_ALIGN.LEFT
                            run.text = str(cell_value)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 0, 0) if row_idx > 0 else RGBColor(255, 255, 255)

    return prs

# Streamlit app
st.title("File Upload and PPT Generation")

# User input for footnote text and year
footnote_text = st.text_input("Enter the footnote text (e.g., month):", "June")
year = st.text_input("Enter the year:", "2024")

uploaded_file = st.file_uploader("Choose a file", type=["xlsx"])

if uploaded_file is not None:
    try:
        df = pd.read_excel(uploaded_file)
        
        st.write("Data uploaded successfully:")
        st.write(df.head())

        prs = create_presentation(df, footnote_text, year)

        pptx_file = f"site_traffic_report_{footnote_text}.pptx"
        with BytesIO() as output:
            prs.save(output)
            output.seek(0)
            st.download_button("Download PPT", output, file_name=pptx_file)

    except Exception as e:
        st.error(f"An error occurred: {e}")
